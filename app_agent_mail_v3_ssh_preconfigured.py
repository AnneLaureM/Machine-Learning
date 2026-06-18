"""
Agent Mail IA — Mailbox Scanner v3.0
====================================

Streamlit app that scans a mailbox through IMAP, summarizes emails with Ollama,
extracts supported attachments, proposes AI-generated replies, and sends replies
through a configured SMTP server.

This script follows the same privacy-first philosophy as the Thunderbird plugin:
- local FastAPI backend is not required here: the Streamlit app talks directly to Ollama;
- Ollama can run locally or behind a secure SSH tunnel;
- email sending uses the SMTP server configured in the app;
- no third-party AI API is used.

Run:
    streamlit run app_agent_mail_v3_ssh_preconfigured.py

Default remote Ollama access:
    ssh -p 22003 sysadmin@195.221.220.18
    local tunnel: 127.0.0.1:11435 -> remote 127.0.0.1:11434

No Ollama port needs to be open publicly.

Recommended dependencies:
    pip install streamlit pandas plotly httpx pypdf python-docx openpyxl python-pptx odfpy beautifulsoup4
"""

from __future__ import annotations

import email
import html
import imaplib
import io
import json
import os
import re
import smtplib
import socket
import subprocess
import time
from dataclasses import dataclass, asdict
from datetime import datetime
from email.header import decode_header, make_header
from email.message import EmailMessage, Message
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.utils import parseaddr, getaddresses
from typing import Any, Dict, Iterable, List, Optional, Tuple

import pandas as pd
import streamlit as st


# =============================================================================
# PAGE CONFIG
# =============================================================================

st.set_page_config(
    page_title="Agent Mail IA — Mailbox Scanner",
    page_icon="📬",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown(
    """
<style>
.block-container {
    padding-top: 1.2rem;
    padding-bottom: 2rem;
}
.agent-card {
    background: #ffffff;
    border: 1px solid #e8eef5;
    border-radius: 14px;
    padding: 16px 18px;
    margin: 10px 0;
    box-shadow: 0 1px 6px rgba(20, 30, 50, 0.06);
}
.agent-card-critical { border-left: 6px solid #d90429; }
.agent-card-high     { border-left: 6px solid #f77f00; }
.agent-card-normal   { border-left: 6px solid #277da1; }
.agent-card-low      { border-left: 6px solid #8d99ae; }
.agent-muted {
    color: #64748b;
    font-size: 0.9rem;
}
.agent-good {
    color: #0f766e;
    font-weight: 600;
}
.agent-warn {
    color: #b45309;
    font-weight: 600;
}
.agent-error {
    color: #b91c1c;
    font-weight: 600;
}
.chat-user {
    background: #eef6ff;
    border-radius: 14px 14px 3px 14px;
    padding: 10px 14px;
    margin: 7px 0 7px 40px;
}
.chat-ai {
    background: #f0fdf4;
    border-radius: 14px 14px 14px 3px;
    padding: 10px 14px;
    margin: 7px 40px 7px 0;
    border-left: 4px solid #22c55e;
}
.log-box {
    background: #111827;
    color: #c7f9cc;
    border-radius: 10px;
    padding: 12px;
    font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, monospace;
    font-size: 0.82rem;
    height: 320px;
    overflow-y: auto;
}
</style>
""",
    unsafe_allow_html=True,
)


# =============================================================================
# SESSION STATE
# =============================================================================

DEFAULTS = {
    "mails": [],
    "logs": [],
    "conversation_by_uid": {},
    "selected_uid": None,
    "triage_done": False,
    "ollama_ok": False,
    "ollama_models": [],
    "ssh_tunnel_proc": None,
    "ssh_tunnel_active": False,
    "sent_history": [],
}

for key, value in DEFAULTS.items():
    if key not in st.session_state:
        st.session_state[key] = value


URGENCY_ORDER = {"critical": 0, "high": 1, "normal": 2, "low": 3}
URGENCY_ICON = {"critical": "🔴", "high": "🟠", "normal": "🟡", "low": "⚪"}
URGENCY_CSS = {"critical": "critical", "high": "high", "normal": "normal", "low": "low"}

DEFAULT_MODEL = "qwen2.5:7b-instruct-q4_K_M"


# =============================================================================
# LOGGING
# =============================================================================

def log(message: str, level: str = "INFO") -> None:
    icon = {
        "INFO": "ℹ️",
        "OK": "✅",
        "WARN": "⚠️",
        "ERR": "❌",
        "LLM": "🧠",
        "MAIL": "📧",
        "SMTP": "📤",
        "SSH": "🔐",
        "PJ": "📎",
    }.get(level, "•")
    ts = datetime.now().strftime("%H:%M:%S")
    st.session_state.logs.append(f"[{ts}] {icon} {message}")


# =============================================================================
# DATA MODELS
# =============================================================================

@dataclass
class AttachmentInfo:
    filename: str
    content_type: str
    size_kb: float
    data: bytes
    extracted_text: str = ""
    extraction_status: str = "not_parsed"


@dataclass
class MailItem:
    uid: str
    message_id: str
    folder: str
    sender: str
    to: str
    cc: str
    subject: str
    date: str
    body_text: str
    body_html: str
    attachments: List[AttachmentInfo]
    analysis: Optional[Dict[str, Any]] = None
    final_draft: str = ""
    sent: bool = False


# =============================================================================
# LOW-LEVEL HELPERS
# =============================================================================

def decode_mime(value: Optional[str]) -> str:
    if not value:
        return ""
    try:
        return str(make_header(decode_header(value))).strip()
    except Exception:
        return value


def normalize_text(text: str, max_chars: Optional[int] = None) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    text = text.strip()
    if max_chars and len(text) > max_chars:
        return text[:max_chars] + "\n\n[... truncated ...]"
    return text


def strip_html(html_text: str) -> str:
    if not html_text:
        return ""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_text, "html.parser")
        return normalize_text(soup.get_text("\n"))
    except Exception:
        text = re.sub(r"<br\s*/?>", "\n", html_text, flags=re.I)
        text = re.sub(r"<[^>]+>", " ", text)
        return normalize_text(html.unescape(text))


def extract_email_address(value: str) -> str:
    return parseaddr(value or "")[1] or value


def safe_subject_reply(subject: str) -> str:
    subject = subject.strip() or "(no subject)"
    return subject if subject.lower().startswith("re:") else f"Re: {subject}"


def parse_addresses(value: str) -> str:
    pairs = getaddresses([value or ""])
    addresses = [addr for _name, addr in pairs if addr]
    return ", ".join(addresses)


# =============================================================================
# OLLAMA CLIENT
# =============================================================================

def check_ollama(url: str) -> Tuple[bool, List[str], str]:
    try:
        import httpx
        with httpx.Client(timeout=5) as client:
            response = client.get(f"{url.rstrip('/')}/api/tags")
            response.raise_for_status()
            data = response.json()
            models = [m.get("name", "") for m in data.get("models", []) if m.get("name")]
        return True, models, "ok"
    except Exception as exc:
        return False, [], str(exc)


def choose_model(models: List[str], preferred: str) -> str:
    if preferred in models:
        return preferred
    if models:
        return models[0]
    return preferred


def ollama_chat(
    messages: List[Dict[str, str]],
    model: str,
    url: str,
    temperature: float = 0.3,
    num_ctx: int = 16384,
) -> str:
    import httpx

    payload = {
        "model": model,
        "messages": messages,
        "stream": False,
        "options": {
            "temperature": temperature,
            "num_ctx": num_ctx,
        },
    }

    try:
        with httpx.Client(timeout=300) as client:
            response = client.post(f"{url.rstrip('/')}/api/chat", json=payload)
            response.raise_for_status()
            data = response.json()
            return data["message"]["content"].strip()
    except httpx.HTTPStatusError as exc:
        detail = exc.response.text[:1000] if exc.response is not None else ""
        raise RuntimeError(f"Ollama HTTP error: {exc}. Details: {detail}") from exc
    except Exception as exc:
        raise RuntimeError(f"Ollama request failed: {exc}") from exc


# =============================================================================
# SSH TUNNEL
# =============================================================================

def open_ssh_tunnel(
    ssh_host: str,
    ssh_port: int,
    ssh_user: str,
    remote_ollama_port: int,
    local_port: int,
    ssh_key: str = "",
) -> Tuple[bool, str, Optional[subprocess.Popen]]:
    import shutil

    ssh_bin = shutil.which("ssh")
    if not ssh_bin:
        return False, "ssh command not found", None

    cmd = [
        ssh_bin,
        "-N",
        "-L", f"{local_port}:127.0.0.1:{remote_ollama_port}",
        "-p", str(ssh_port),
        "-o", "ExitOnForwardFailure=yes",
        "-o", "ServerAliveInterval=30",
        "-o", "ServerAliveCountMax=3",
        "-o", "ConnectTimeout=10",
    ]

    if ssh_key and os.path.isfile(os.path.expanduser(ssh_key)):
        cmd += ["-i", os.path.expanduser(ssh_key)]

    cmd.append(f"{ssh_user}@{ssh_host}")

    try:
        proc = subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE)
        time.sleep(1.5)
        if proc.poll() is not None:
            err = proc.stderr.read().decode(errors="replace")
            return False, err[:800], None
        return True, f"Tunnel active: localhost:{local_port} → {ssh_host}:127.0.0.1:{remote_ollama_port}", proc
    except Exception as exc:
        return False, str(exc), None


def close_ssh_tunnel() -> None:
    proc = st.session_state.get("ssh_tunnel_proc")
    if proc:
        try:
            proc.terminate()
        except Exception:
            pass
    st.session_state.ssh_tunnel_proc = None
    st.session_state.ssh_tunnel_active = False


# =============================================================================
# IMAP READER
# =============================================================================

class MailboxReader:
    def __init__(self, imap_host: str, imap_port: int, user: str, password: str, folder: str = "INBOX"):
        self.imap_host = imap_host
        self.imap_port = imap_port
        self.user = user
        self.password = password
        self.folder = folder
        self.conn: Optional[imaplib.IMAP4_SSL] = None

    def connect(self) -> Tuple[bool, str]:
        try:
            self.conn = imaplib.IMAP4_SSL(self.imap_host, self.imap_port)
            self.conn.login(self.user, self.password)
            return True, "connected"
        except Exception as exc:
            return False, str(exc)

    def disconnect(self) -> None:
        try:
            if self.conn:
                self.conn.logout()
        except Exception:
            pass

    def search_uids(self, unread_only: bool, since: str = "") -> List[bytes]:
        if not self.conn:
            return []
        self.conn.select(self.folder)
        criteria = ["UNSEEN"] if unread_only else ["ALL"]
        if since:
            criteria += ["SINCE", since]
        status, data = self.conn.search(None, *criteria)
        if status != "OK" or not data:
            return []
        return data[0].split()

    def fetch(self, n: int = 20, unread_only: bool = False, since: str = "") -> List[MailItem]:
        if not self.conn:
            return []

        uids = self.search_uids(unread_only=unread_only, since=since)
        uids = uids[-n:][::-1]

        mails: List[MailItem] = []
        for uid in uids:
            try:
                status, data = self.conn.fetch(uid, "(RFC822)")
                if status != "OK" or not data or not isinstance(data[0], tuple):
                    continue
                msg = email.message_from_bytes(data[0][1])
                mails.append(self.parse_message(uid.decode(), msg))
            except Exception as exc:
                log(f"Unable to parse mail {uid!r}: {exc}", "ERR")

        return mails

    def parse_message(self, uid: str, msg: Message) -> MailItem:
        body_text, body_html = self.extract_bodies(msg)
        attachments = self.extract_attachments(msg)

        return MailItem(
            uid=uid,
            message_id=decode_mime(msg.get("Message-ID", "")),
            folder=self.folder,
            sender=decode_mime(msg.get("From", "")),
            to=decode_mime(msg.get("To", "")),
            cc=decode_mime(msg.get("Cc", "")),
            subject=decode_mime(msg.get("Subject", "")) or "(no subject)",
            date=decode_mime(msg.get("Date", "")),
            body_text=normalize_text(body_text or strip_html(body_html), max_chars=12000),
            body_html=body_html,
            attachments=attachments,
        )

    def extract_bodies(self, msg: Message) -> Tuple[str, str]:
        plain_parts: List[str] = []
        html_parts: List[str] = []

        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                disposition = str(part.get("Content-Disposition", "")).lower()
                if "attachment" in disposition:
                    continue
                payload = part.get_payload(decode=True)
                if payload is None:
                    continue
                charset = part.get_content_charset() or "utf-8"
                try:
                    text = payload.decode(charset, errors="replace")
                except Exception:
                    text = payload.decode("utf-8", errors="replace")
                if content_type == "text/plain":
                    plain_parts.append(text)
                elif content_type == "text/html":
                    html_parts.append(text)
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or "utf-8"
                text = payload.decode(charset, errors="replace")
                if msg.get_content_type() == "text/html":
                    html_parts.append(text)
                else:
                    plain_parts.append(text)

        return normalize_text("\n".join(plain_parts)), "\n".join(html_parts)

    def extract_attachments(self, msg: Message) -> List[AttachmentInfo]:
        attachments: List[AttachmentInfo] = []
        for part in msg.walk():
            filename = decode_mime(part.get_filename())
            if not filename:
                continue
            payload = part.get_payload(decode=True) or b""
            attachments.append(
                AttachmentInfo(
                    filename=filename,
                    content_type=part.get_content_type(),
                    size_kb=round(len(payload) / 1024, 1),
                    data=payload,
                )
            )
        return attachments


# =============================================================================
# ATTACHMENT EXTRACTION
# =============================================================================

class AttachmentParser:
    def __init__(self, max_chars_per_attachment: int = 6000):
        self.max_chars_per_attachment = max_chars_per_attachment

    def parse(self, att: AttachmentInfo) -> AttachmentInfo:
        filename = att.filename.lower()
        content_type = (att.content_type or "").lower()

        try:
            if filename.endswith((".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm")) or content_type.startswith("text/"):
                text = self._parse_text_like(att)
            elif filename.endswith(".pdf") or "pdf" in content_type:
                text = self._parse_pdf(att)
            elif filename.endswith(".docx"):
                text = self._parse_docx(att)
            elif filename.endswith((".xlsx", ".xlsm", ".xls")):
                text = self._parse_excel(att)
            elif filename.endswith(".pptx"):
                text = self._parse_pptx(att)
            elif filename.endswith(".odt"):
                text = self._parse_odt(att)
            elif content_type.startswith("image/") or filename.endswith((".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff")):
                text = f"[Image attachment: {att.filename}. OCR is not implemented in this script.]"
                att.extraction_status = "unsupported_ocr_needed"
                att.extracted_text = text
                return att
            else:
                text = f"[Unsupported attachment type: {att.filename} / {att.content_type}]"
                att.extraction_status = "unsupported"
                att.extracted_text = text
                return att

            att.extracted_text = normalize_text(text, self.max_chars_per_attachment)
            att.extraction_status = "ok" if att.extracted_text.strip() else "empty"
            return att

        except Exception as exc:
            att.extracted_text = f"[Attachment extraction error for {att.filename}: {exc}]"
            att.extraction_status = "error"
            return att

    def _parse_text_like(self, att: AttachmentInfo) -> str:
        text = att.data.decode("utf-8", errors="replace")
        if att.filename.lower().endswith(".csv"):
            try:
                df = pd.read_csv(io.StringIO(text))
                return (
                    f"[CSV: {att.filename}]\n"
                    f"Rows: {len(df)}\n"
                    f"Columns: {list(df.columns)}\n\n"
                    f"Preview:\n{df.head(10).to_string(index=False)}"
                )
            except Exception:
                pass
        if att.filename.lower().endswith((".html", ".htm")):
            return strip_html(text)
        return f"[Text: {att.filename}]\n{text}"

    def _parse_pdf(self, att: AttachmentInfo) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise RuntimeError("Please install pypdf: pip install pypdf")

        reader = PdfReader(io.BytesIO(att.data))
        pages = []
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text() or ""
            except Exception:
                page_text = ""
            pages.append(f"\n--- Page {i + 1} ---\n{page_text}")
        return f"[PDF: {att.filename} | {len(reader.pages)} pages]\n" + "\n".join(pages)

    def _parse_docx(self, att: AttachmentInfo) -> str:
        try:
            import docx
        except ImportError:
            raise RuntimeError("Please install python-docx: pip install python-docx")

        doc = docx.Document(io.BytesIO(att.data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows[:20]:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                tables.append(f"\n[Table {table_idx + 1}]\n" + "\n".join(rows))
        return f"[DOCX: {att.filename}]\n" + "\n".join(paragraphs + tables)

    def _parse_excel(self, att: AttachmentInfo) -> str:
        try:
            xls = pd.ExcelFile(io.BytesIO(att.data))
        except Exception as exc:
            raise RuntimeError(f"Please install openpyxl for xlsx files. Details: {exc}")

        parts = [f"[Excel: {att.filename}]"]
        for sheet in xls.sheet_names[:5]:
            df = pd.read_excel(xls, sheet_name=sheet)
            parts.append(
                f"\n--- Sheet: {sheet} | rows={len(df)} cols={len(df.columns)} ---\n"
                f"Columns: {list(df.columns)}\n"
                f"{df.head(12).to_string(index=False)}"
            )
        return "\n".join(parts)

    def _parse_pptx(self, att: AttachmentInfo) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("Please install python-pptx: pip install python-pptx")

        prs = Presentation(io.BytesIO(att.data))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"\n--- Slide {i + 1} ---\n" + "\n".join(texts))
        return f"[PPTX: {att.filename} | {len(prs.slides)} slides]\n" + "\n".join(slides)

    def _parse_odt(self, att: AttachmentInfo) -> str:
        try:
            from odf.opendocument import load
            from odf import text as odf_text
        except ImportError:
            raise RuntimeError("Please install odfpy: pip install odfpy")

        doc = load(io.BytesIO(att.data))
        paragraphs = []
        for p in doc.getElementsByType(odf_text.P):
            paragraphs.append("".join(node.data for node in p.childNodes if hasattr(node, "data")))
        return f"[ODT: {att.filename}]\n" + "\n".join(paragraphs)


PARSER = AttachmentParser()


def parse_all_attachments(mail: MailItem) -> MailItem:
    parsed = []
    for att in mail.attachments:
        parsed.append(PARSER.parse(att))
    mail.attachments = parsed
    return mail


def attachment_context(mail: MailItem, max_total_chars: int = 18000) -> str:
    parts = []
    for att in mail.attachments:
        if att.extracted_text:
            parts.append(
                f"--- {att.filename} ({att.content_type}, {att.size_kb} KB, status={att.extraction_status}) ---\n"
                f"{att.extracted_text}"
            )
    text = "\n\n".join(parts)
    return normalize_text(text, max_total_chars)


def attachment_diagnostics(mail: MailItem) -> Tuple[int, int, List[str]]:
    total = len(mail.attachments)
    ok = sum(1 for a in mail.attachments if a.extraction_status == "ok")
    lines = []
    for att in mail.attachments:
        if att.extraction_status == "ok":
            lines.append(f"✓ {att.filename} — {len(att.extracted_text)} chars extracted")
        elif att.extraction_status == "not_parsed":
            lines.append(f"• {att.filename} — not parsed yet")
        else:
            lines.append(f"⚠ {att.filename} — {att.extraction_status}")
    return ok, total, lines


# =============================================================================
# LLM PROMPTS
# =============================================================================

def build_mail_context(mail: MailItem, include_attachments: bool = True) -> str:
    att_names = ", ".join(a.filename for a in mail.attachments) if mail.attachments else "none"
    att_text = attachment_context(mail) if include_attachments else ""

    return f"""EMAIL CONTEXT
From: {mail.sender}
To: {mail.to}
Cc: {mail.cc}
Subject: {mail.subject}
Date: {mail.date}
Attachments: {att_names}

EMAIL BODY
{normalize_text(mail.body_text, 10000)}

ATTACHMENT CONTENT
{att_text if att_text else "[No extracted attachment content available.]"}
"""


def analyze_mail(mail: MailItem, model: str, ollama_url: str) -> Dict[str, Any]:
    schema = {
        "urgency": "critical|high|normal|low",
        "category": "action_required|question|information|administrative|commercial|spam|other",
        "summary": "short factual summary in English or French",
        "key_points": ["point 1", "point 2"],
        "attachment_summary": "summary of useful attachment content, or empty string",
        "suggested_action": "what the user should do next",
        "reply_needed": True,
        "deadline": "immediate|today|this_week|none|unknown",
        "confidence": 0.0,
    }

    prompt = f"""Analyze the email below.

Return ONLY valid JSON matching this schema:
{json.dumps(schema, ensure_ascii=False)}

Rules:
- Use French unless the email is clearly in English.
- If attachments contain useful information, summarize them explicitly.
- Do not invent facts.
- If no reply is needed, set reply_needed=false.

{build_mail_context(mail, include_attachments=True)}
"""

    raw = ollama_chat(
        [
            {"role": "system", "content": "You are a precise email triage assistant. You return strict JSON only."},
            {"role": "user", "content": prompt},
        ],
        model=model,
        url=ollama_url,
        temperature=0.05,
    )

    raw = raw.replace("```json", "").replace("```", "").strip()
    start, end = raw.find("{"), raw.rfind("}") + 1
    if start >= 0 and end > start:
        raw = raw[start:end]

    try:
        data = json.loads(raw)
        if data.get("urgency") not in URGENCY_ORDER:
            data["urgency"] = "normal"
        data["confidence"] = float(data.get("confidence", 0.5))
        return data
    except Exception:
        return {
            "urgency": "normal",
            "category": "information",
            "summary": mail.subject,
            "key_points": [],
            "attachment_summary": "",
            "suggested_action": "Review manually",
            "reply_needed": True,
            "deadline": "unknown",
            "confidence": 0.5,
        }


def build_reply_system_prompt(mail: MailItem) -> str:
    analysis = mail.analysis or {}
    return f"""You are a professional email assistant.

You draft replies that are ready to review and send through the user's configured email server.

Rules:
- Write in the same language as the incoming email, unless the user requests otherwise.
- Be concise, professional, and natural.
- Use the attachment content when relevant or when the user asks for it.
- Do not invent information that is not present in the email or attachments.
- Do not include fake signatures, job titles, or names.
- Do not add commentary outside the draft unless explicitly asked.
- For a draft reply, start with an appropriate greeting and end with "Cordialement," or the corresponding closing in the email language.
- The user will send the email through the configured SMTP server.

Current triage:
Urgency: {analysis.get("urgency", "unknown")}
Category: {analysis.get("category", "unknown")}
Summary: {analysis.get("summary", "")}
Suggested action: {analysis.get("suggested_action", "")}
"""


def chat_about_mail(
    mail: MailItem,
    user_message: str,
    model: str,
    ollama_url: str,
    temperature: float,
) -> str:
    uid = mail.uid
    st.session_state.conversation_by_uid.setdefault(uid, [])

    conversation = st.session_state.conversation_by_uid[uid]
    conversation.append({"role": "user", "content": user_message})

    messages = [
        {"role": "system", "content": build_reply_system_prompt(mail)},
        {"role": "user", "content": build_mail_context(mail, include_attachments=True)},
    ] + conversation

    response = ollama_chat(messages, model=model, url=ollama_url, temperature=temperature)
    conversation.append({"role": "assistant", "content": response})
    st.session_state.conversation_by_uid[uid] = conversation
    return response


# =============================================================================
# SMTP SENDER
# =============================================================================

def send_email_smtp(
    smtp_host: str,
    smtp_port: int,
    smtp_user: str,
    smtp_password: str,
    sender_email: str,
    to: str,
    subject: str,
    body: str,
    cc: str = "",
    bcc: str = "",
    use_ssl: bool = True,
    use_starttls: bool = False,
    reply_to_message_id: str = "",
) -> Tuple[bool, str]:
    try:
        msg = MIMEMultipart("alternative")
        msg["From"] = sender_email or smtp_user
        msg["To"] = to
        if cc:
            msg["Cc"] = cc
        msg["Subject"] = subject
        if reply_to_message_id:
            msg["In-Reply-To"] = reply_to_message_id
            msg["References"] = reply_to_message_id

        msg.attach(MIMEText(body, "plain", "utf-8"))

        recipients = []
        for field in [to, cc, bcc]:
            recipients.extend([addr for _name, addr in getaddresses([field]) if addr])

        if use_ssl:
            server = smtplib.SMTP_SSL(smtp_host, smtp_port, timeout=60)
        else:
            server = smtplib.SMTP(smtp_host, smtp_port, timeout=60)

        with server as smtp:
            if use_starttls and not use_ssl:
                smtp.starttls()
            if smtp_user:
                smtp.login(smtp_user, smtp_password)
            smtp.send_message(msg, from_addr=sender_email or smtp_user, to_addrs=recipients)

        return True, "sent"

    except Exception as exc:
        return False, str(exc)


# =============================================================================
# SIDEBAR CONFIG
# =============================================================================

with st.sidebar:
    st.title("⚙️ Configuration")

    st.markdown("### 📥 Mailbox / IMAP")
    imap_preset = st.selectbox("Provider preset", ["Gmail", "Outlook / Office365", "Custom"], index=0)

    if imap_preset == "Gmail":
        default_imap_host, default_imap_port = "imap.gmail.com", 993
        default_smtp_host, default_smtp_port, default_ssl = "smtp.gmail.com", 465, True
    elif imap_preset == "Outlook / Office365":
        default_imap_host, default_imap_port = "outlook.office365.com", 993
        default_smtp_host, default_smtp_port, default_ssl = "smtp.office365.com", 587, False
    else:
        default_imap_host, default_imap_port = "", 993
        default_smtp_host, default_smtp_port, default_ssl = "", 587, False

    email_user = st.text_input("Mailbox username / email", value="", placeholder="name@example.com")
    email_password = st.text_input("Mailbox password / app password", value="", type="password")

    imap_host = st.text_input("IMAP host", value=default_imap_host)
    imap_port = st.number_input("IMAP port", value=default_imap_port, step=1)
    imap_folder = st.text_input("IMAP folder", value="INBOX")

    st.markdown("### 📤 Sending / SMTP")
    st.caption("A Streamlit app cannot read Thunderbird's native SMTP settings. Configure the same outgoing server here.")
    smtp_host = st.text_input("SMTP host", value=default_smtp_host)
    smtp_port = st.number_input("SMTP port", value=default_smtp_port, step=1)
    smtp_user = st.text_input("SMTP username", value=email_user)
    smtp_password = st.text_input("SMTP password / app password", value="", type="password")
    sender_email = st.text_input("From address", value=email_user)
    smtp_ssl = st.checkbox("SMTP SSL", value=default_ssl)
    smtp_starttls = st.checkbox("SMTP STARTTLS", value=(imap_preset == "Outlook / Office365"))

    st.markdown("### 🦙 Ollama")
    inference_mode = st.radio("Inference mode", ["Local", "Remote via SSH"], horizontal=True)

    if inference_mode == "Local":
        ollama_url = st.text_input("Ollama URL", value="http://127.0.0.1:11434")
        if st.session_state.ssh_tunnel_active:
            close_ssh_tunnel()
            log("SSH tunnel closed because local mode was selected.", "SSH")
    else:
        ssh_host = st.text_input("SSH host", value="195.221.220.18")
        ssh_port = st.number_input("SSH port", value=22003, step=1)
        ssh_user = st.text_input("SSH user", value="sysadmin")
        ssh_key = st.text_input("SSH private key path", value=os.path.expanduser("~/.ssh/id_rsa"))
        remote_ollama_port = st.number_input("Remote Ollama port", value=11434, step=1)
        local_tunnel_port = st.number_input("Local tunnel port", value=11435, step=1)
        # Ollama is never reached directly over the network.
        # The app only calls this local forwarded port.
        # SSH forwards it to 127.0.0.1:11434 on the remote server.
        ollama_url = f"http://127.0.0.1:{int(local_tunnel_port)}"
        st.caption("Ollama remote access is done only through SSH: no public Ollama port is required.")

        col_open, col_close = st.columns(2)
        with col_open:
            if st.button("🔐 Open SSH tunnel", disabled=st.session_state.ssh_tunnel_active):
                ok, msg, proc = open_ssh_tunnel(
                    ssh_host=ssh_host,
                    ssh_port=int(ssh_port),
                    ssh_user=ssh_user,
                    remote_ollama_port=int(remote_ollama_port),
                    local_port=int(local_tunnel_port),
                    ssh_key=ssh_key,
                )
                if ok:
                    st.session_state.ssh_tunnel_proc = proc
                    st.session_state.ssh_tunnel_active = True
                    log(msg, "SSH")
                    st.rerun()
                else:
                    st.error(msg)
                    log(f"SSH tunnel failed: {msg}", "ERR")
        with col_close:
            if st.button("✂️ Close tunnel", disabled=not st.session_state.ssh_tunnel_active):
                close_ssh_tunnel()
                log("SSH tunnel closed.", "SSH")
                st.rerun()

        if st.session_state.ssh_tunnel_active:
            proc = st.session_state.get("ssh_tunnel_proc")
            if proc and proc.poll() is None:
                st.success(f"SSH tunnel active → {ollama_url}")
            else:
                st.warning("SSH tunnel process is no longer active.")
                st.session_state.ssh_tunnel_active = False
        else:
            st.info("SSH tunnel not active.")

        with st.expander("Equivalent SSH command"):
            st.code(
                f"ssh -N -L {int(local_tunnel_port)}:127.0.0.1:{int(remote_ollama_port)} "
                f"-p {int(ssh_port)} -i {ssh_key} {ssh_user}@{ssh_host}",
                language="bash",
            )

    ok_ollama, ollama_models, ollama_error = check_ollama(ollama_url)
    st.session_state.ollama_ok = ok_ollama
    st.session_state.ollama_models = ollama_models

    if ok_ollama:
        st.success(f"Ollama reachable — {len(ollama_models)} model(s)")
        selected_model = st.selectbox(
            "Model",
            options=ollama_models,
            index=ollama_models.index(DEFAULT_MODEL) if DEFAULT_MODEL in ollama_models else 0,
        )
    else:
        st.error(f"Ollama unreachable: {ollama_error}")
        selected_model = st.text_input("Model name", value=DEFAULT_MODEL)

    temperature = st.slider("Reply temperature", 0.0, 1.0, 0.35, 0.05)

    st.markdown("### 📬 Scan options")
    nb_mails = st.slider("Number of emails", 5, 100, 20)
    unread_only = st.checkbox("Unread only", value=False)
    scan_since = st.text_input("IMAP SINCE filter", value="", placeholder="01-Jun-2026")
    auto_parse_attachments = st.checkbox("Extract attachments automatically", value=True)
    auto_triage = st.checkbox("Analyze emails automatically after scan", value=True)

    st.divider()
    if st.button("🧹 Reset session"):
        keep_keys = []
        for key in list(st.session_state.keys()):
            if key not in keep_keys:
                del st.session_state[key]
        st.rerun()


# =============================================================================
# MAIN UI
# =============================================================================

st.title("📬 Agent Mail IA — Mailbox Scanner")
st.caption("Mailbox scan · LLM summaries · attachment-aware replies · SMTP sending · local/SSH Ollama")

tab_scan, tab_triage, tab_reply, tab_logs = st.tabs(["📥 Scan", "🧠 Summaries", "✉️ Reply", "🖥 Logs"])


# -----------------------------------------------------------------------------
# TAB SCAN
# -----------------------------------------------------------------------------

with tab_scan:
    st.subheader("Scan mailbox")

    col_status_1, col_status_2, col_status_3 = st.columns(3)
    col_status_1.metric("Ollama", "OK" if st.session_state.ollama_ok else "Unavailable")
    col_status_2.metric("Loaded emails", len(st.session_state.mails))
    col_status_3.metric("Conversations", len(st.session_state.conversation_by_uid))

    if st.button("📨 Scan mailbox", type="primary", disabled=not (email_user and email_password and imap_host)):
        with st.spinner("Connecting to mailbox..."):
            reader = MailboxReader(
                imap_host=imap_host,
                imap_port=int(imap_port),
                user=email_user,
                password=email_password,
                folder=imap_folder,
            )
            ok, msg = reader.connect()
            if not ok:
                st.error(f"IMAP error: {msg}")
                log(f"IMAP connection error: {msg}", "ERR")
            else:
                try:
                    mails = reader.fetch(n=nb_mails, unread_only=unread_only, since=scan_since.strip())
                    if auto_parse_attachments:
                        for mail in mails:
                            if mail.attachments:
                                parse_all_attachments(mail)
                    st.session_state.mails = mails
                    st.session_state.selected_uid = mails[0].uid if mails else None
                    st.session_state.triage_done = False
                    log(f"{len(mails)} email(s) loaded from {imap_folder}.", "MAIL")
                    st.success(f"{len(mails)} email(s) loaded.")
                finally:
                    reader.disconnect()

        if auto_triage and st.session_state.mails and st.session_state.ollama_ok:
            with st.spinner("Analyzing emails with Ollama..."):
                progress = st.progress(0)
                for i, mail in enumerate(st.session_state.mails):
                    try:
                        mail.analysis = analyze_mail(mail, selected_model, ollama_url)
                        log(f"Analyzed: {mail.subject[:60]}", "LLM")
                    except Exception as exc:
                        log(f"Analysis failed for {mail.subject[:40]}: {exc}", "ERR")
                    progress.progress((i + 1) / len(st.session_state.mails))
                st.session_state.triage_done = True
                st.success("Analysis complete.")
                st.rerun()

    if st.session_state.mails:
        st.markdown("### Loaded emails")
        sorted_mails = sorted(
            st.session_state.mails,
            key=lambda m: URGENCY_ORDER.get((m.analysis or {}).get("urgency", "normal"), 2),
        )

        for mail in sorted_mails:
            analysis = mail.analysis or {}
            urgency = analysis.get("urgency", "normal")
            icon = URGENCY_ICON.get(urgency, "📧")
            css = URGENCY_CSS.get(urgency, "normal")
            ok_att, total_att, att_lines = attachment_diagnostics(mail)

            with st.expander(f"{icon} {mail.subject[:90]} — {mail.sender[:50]}"):
                st.markdown(
                    f"""
<div class="agent-card agent-card-{css}">
<b>From:</b> {html.escape(mail.sender)}<br>
<b>To:</b> {html.escape(mail.to)}<br>
<b>Date:</b> {html.escape(mail.date)}<br>
<b>Attachments:</b> {ok_att}/{total_att} usable
</div>
""",
                    unsafe_allow_html=True,
                )

                if analysis:
                    st.markdown(f"**Summary:** {analysis.get('summary', '')}")
                    st.markdown(f"**Suggested action:** {analysis.get('suggested_action', '')}")

                st.text_area("Body preview", mail.body_text[:2500], height=180, disabled=True, key=f"body_{mail.uid}")

                if total_att:
                    st.markdown("**Attachment diagnostics**")
                    for line in att_lines:
                        st.write(line)

                if st.button("Select for reply", key=f"select_{mail.uid}"):
                    st.session_state.selected_uid = mail.uid
                    st.rerun()


# -----------------------------------------------------------------------------
# TAB TRIAGE
# -----------------------------------------------------------------------------

with tab_triage:
    st.subheader("LLM summaries and triage")

    if not st.session_state.mails:
        st.info("Scan the mailbox first.")
    else:
        if st.button("🧠 Analyze / refresh summaries", type="primary", disabled=not st.session_state.ollama_ok):
            with st.spinner("Analyzing emails..."):
                progress = st.progress(0)
                for i, mail in enumerate(st.session_state.mails):
                    if auto_parse_attachments and mail.attachments:
                        parse_all_attachments(mail)
                    mail.analysis = analyze_mail(mail, selected_model, ollama_url)
                    progress.progress((i + 1) / len(st.session_state.mails))
                st.session_state.triage_done = True
                st.rerun()

        rows = []
        for mail in st.session_state.mails:
            a = mail.analysis or {}
            rows.append(
                {
                    "urgency": a.get("urgency", ""),
                    "category": a.get("category", ""),
                    "reply_needed": a.get("reply_needed", ""),
                    "deadline": a.get("deadline", ""),
                    "summary": a.get("summary", ""),
                    "from": mail.sender,
                    "subject": mail.subject,
                    "attachments": len(mail.attachments),
                    "sent": mail.sent,
                }
            )

        if rows:
            df = pd.DataFrame(rows)
            st.dataframe(df, use_container_width=True, hide_index=True)

            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Emails", len(df))
            col2.metric("Critical", int((df["urgency"] == "critical").sum()))
            col3.metric("High", int((df["urgency"] == "high").sum()))
            col4.metric("Replies needed", int((df["reply_needed"] == True).sum()))


# -----------------------------------------------------------------------------
# TAB REPLY
# -----------------------------------------------------------------------------

with tab_reply:
    st.subheader("Draft reply and send through configured SMTP")

    if not st.session_state.mails:
        st.info("Scan the mailbox first.")
    else:
        uid_to_mail = {m.uid: m for m in st.session_state.mails}
        selected_uid = st.session_state.selected_uid or next(iter(uid_to_mail.keys()))
        uid_options = list(uid_to_mail.keys())

        selected_uid = st.selectbox(
            "Email",
            options=uid_options,
            index=uid_options.index(selected_uid) if selected_uid in uid_options else 0,
            format_func=lambda uid: f"{uid_to_mail[uid].subject[:80]} — {uid_to_mail[uid].sender[:40]}",
        )
        st.session_state.selected_uid = selected_uid
        mail = uid_to_mail[selected_uid]

        analysis = mail.analysis or {}
        ok_att, total_att, att_lines = attachment_diagnostics(mail)

        with st.expander("Selected email context", expanded=True):
            st.markdown(f"**From:** {mail.sender}")
            st.markdown(f"**Subject:** {mail.subject}")
            st.markdown(f"**Summary:** {analysis.get('summary', 'Not analyzed yet')}")
            if total_att:
                st.markdown(f"**Attachments:** {ok_att}/{total_att} usable")
                for line in att_lines:
                    st.write(line)

        st.markdown("### Conversation with the LLM")

        conversation = st.session_state.conversation_by_uid.get(mail.uid, [])
        for msg in conversation:
            klass = "chat-user" if msg["role"] == "user" else "chat-ai"
            who = "You" if msg["role"] == "user" else "LLM"
            st.markdown(
                f'<div class="{klass}"><b>{who}</b><br>{html.escape(msg["content"]).replace(chr(10), "<br>")}</div>',
                unsafe_allow_html=True,
            )

        col_shortcuts = st.columns(4)
        shortcut_prompt = None
        with col_shortcuts[0]:
            if st.button("📝 Draft reply"):
                shortcut_prompt = "Generate a professional reply to this email. Use attachments if relevant."
        with col_shortcuts[1]:
            if st.button("📎 Use attachments"):
                shortcut_prompt = "Generate a professional reply that explicitly takes into account the useful information from the attachments."
        with col_shortcuts[2]:
            if st.button("✂️ Shorter"):
                shortcut_prompt = "Rewrite the last draft in a shorter and more direct style."
        with col_shortcuts[3]:
            if st.button("🤝 Warmer"):
                shortcut_prompt = "Rewrite the last draft with a warmer and more collaborative tone."

        user_prompt_key = f"user_prompt_{mail.uid}"
        if shortcut_prompt:
            st.session_state[user_prompt_key] = shortcut_prompt
            st.rerun()

        user_prompt = st.text_area(
            "Instruction",
            key=user_prompt_key,
            height=90,
            placeholder="Example: Draft a reply confirming receipt and mentioning the attached quote.",
        )

        if st.button("🧠 Send to LLM", type="primary", disabled=not (st.session_state.ollama_ok and user_prompt.strip())):
            if auto_parse_attachments and mail.attachments:
                parse_all_attachments(mail)
            with st.spinner("Generating reply..."):
                try:
                    reply = chat_about_mail(mail, user_prompt.strip(), selected_model, ollama_url, temperature)
                    mail.final_draft = reply
                    st.session_state[f"draft_{mail.uid}"] = reply
                    st.rerun()
                except Exception as exc:
                    st.error(f"LLM error: {exc}")
                    log(f"LLM error: {exc}", "ERR")

        st.markdown("### Final draft")

        draft_key = f"draft_{mail.uid}"
        if draft_key not in st.session_state:
            st.session_state[draft_key] = mail.final_draft or ""

        col_to, col_subject = st.columns(2)
        with col_to:
            to_default = parse_addresses(mail.sender)
            to_addr = st.text_input("To", value=to_default, key=f"to_{mail.uid}")
        with col_subject:
            subject = st.text_input("Subject", value=safe_subject_reply(mail.subject), key=f"subject_{mail.uid}")

        cc_addr = st.text_input("Cc", value="", key=f"cc_{mail.uid}")
        bcc_addr = st.text_input("Bcc", value="", key=f"bcc_{mail.uid}")

        final_draft = st.text_area(
            "Reply body",
            key=draft_key,
            height=360,
            placeholder="Generate a draft with the LLM, then edit it here before sending.",
        )

        with st.expander("Sending configuration recap", expanded=True):
            st.markdown(f"**From:** `{sender_email or smtp_user}`")
            st.markdown(f"**SMTP:** `{smtp_host}:{int(smtp_port)}`")
            st.markdown(f"**Security:** {'SSL' if smtp_ssl else 'STARTTLS' if smtp_starttls else 'plain SMTP'}")
            st.markdown(f"**To:** `{to_addr}`")
            st.markdown(f"**Subject:** `{subject}`")

        col_send, col_save = st.columns([2, 1])
        with col_send:
            send_now = st.button("📤 Send reply via configured SMTP", type="primary", disabled=not (smtp_host and smtp_password and final_draft.strip()))
        with col_save:
            st.download_button(
                "💾 Save draft",
                data=final_draft,
                file_name=f"reply_{mail.uid}.txt",
                mime="text/plain",
                use_container_width=True,
            )

        if send_now:
            with st.spinner("Sending email..."):
                ok, msg = send_email_smtp(
                    smtp_host=smtp_host,
                    smtp_port=int(smtp_port),
                    smtp_user=smtp_user,
                    smtp_password=smtp_password,
                    sender_email=sender_email or smtp_user,
                    to=to_addr,
                    subject=subject,
                    body=final_draft,
                    cc=cc_addr,
                    bcc=bcc_addr,
                    use_ssl=smtp_ssl,
                    use_starttls=smtp_starttls,
                    reply_to_message_id=mail.message_id,
                )
            if ok:
                mail.sent = True
                st.session_state.sent_history.append(
                    {
                        "time": datetime.now().isoformat(timespec="seconds"),
                        "to": to_addr,
                        "subject": subject,
                        "smtp": f"{smtp_host}:{int(smtp_port)}",
                    }
                )
                log(f"Email sent to {to_addr}: {subject}", "SMTP")
                st.success("Email sent.")
                st.balloons()
            else:
                st.error(f"SMTP error: {msg}")
                log(f"SMTP error: {msg}", "ERR")

        if st.session_state.sent_history:
            st.markdown("### Recent sends")
            st.dataframe(pd.DataFrame(st.session_state.sent_history), use_container_width=True, hide_index=True)


# -----------------------------------------------------------------------------
# TAB LOGS
# -----------------------------------------------------------------------------

with tab_logs:
    st.subheader("Logs and diagnostics")

    col1, col2, col3 = st.columns(3)
    col1.metric("Logs", len(st.session_state.logs))
    col2.metric("Ollama", "OK" if st.session_state.ollama_ok else "Unavailable")
    col3.metric("Models", len(st.session_state.ollama_models))

    if st.button("Refresh"):
        st.rerun()

    logs = "\n".join(st.session_state.logs[-150:])
    st.markdown(f'<div class="log-box"><pre>{html.escape(logs)}</pre></div>', unsafe_allow_html=True)

    st.markdown("### Ollama models")
    if st.session_state.ollama_models:
        for model_name in st.session_state.ollama_models:
            st.write(f"🦙 {model_name}")
    else:
        st.info("No model listed.")

    st.download_button(
        "Download logs",
        data="\n".join(st.session_state.logs),
        file_name=f"agent_mail_logs_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
        mime="text/plain",
    )
